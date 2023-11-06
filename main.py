# Módulos externos:
# pip install python-pptx
# pip install pyinstaller

# Comando para criar executável, esse comando adiciona os arquivos "favicon.ico" e o "default.pptx do python-pptx:
# pyinstaller --add-data "venv\Lib\site-packages\pptx\templates\default.pptx:." --add-data "favicon.ico:." --icon=favicon.ico --noconsole --onefile main.py

from pptx import Presentation
import tkinter as tk
from tkinter import messagebox
import re
from pptx import Presentation
import os

# Diretório do script atual
script_dir = os.path.dirname(os.path.abspath(__file__))

# Caminho relativo para o arquivo 'default.pptx'
rel_path = "default.pptx"
icone = "favicon.ico"

# Caminho absoluto para o 'default.pptx' usando o caminho relativo
abs_file_path = os.path.join(script_dir, rel_path)
abs_file_path_icone = os.path.join(script_dir, icone)

def clique_botao():
    letra = caixa_texto.get("1.0", "end-1c")
    letraLinhas = letra.split("\n")

    if os.path.exists(abs_file_path):
        root = Presentation(abs_file_path)
    else:
        root = Presentation()

    slideTitulo = root.slide_layouts[0]
    titulo = root.slides.add_slide(slideTitulo)
    titulo.shapes.title.text = letraLinhas[0].strip().upper()
    del letraLinhas[0]

    letraMusica = root.slide_layouts[2]
    for linha in letraLinhas:
        linha = linha.strip()
        if linha:
            slide = root.slides.add_slide(letraMusica)
            slide.shapes.title.text = linha.upper()

    nomeArquivo = re.sub(u'[^a-zA-Z0-9áéíóúÁÉÍÓÚâêîôÂÊÎÔãõÃÕçÇ: ]', '', titulo.shapes.title.text)
    root.save(f'{nomeArquivo}.pptx')

    messagebox.showinfo("Criador de PowerPoint", "Arquivo PowerPoint criado")

def colar_texto():
    texto_copiado = janela.clipboard_get()
    caixa_texto.insert("insert", texto_copiado)

def copiar_texto():
    texto_selecionado = caixa_texto.selection_get()  # Obtém o texto selecionado na caixa de texto
    janela.clipboard_clear()  # Limpa a área de transferência
    janela.clipboard_append(texto_selecionado)  # Copia o texto selecionado para a área de transferência

def abrir_menu(event):
    menu_contextual.post(event.x_root, event.y_root)

# Criando a janela principal
janela = tk.Tk()
janela.title("Criador de PowerPoint de letra de música")
janela.iconbitmap(abs_file_path_icone)

# Rótulo para o texto fixo
legenda = tk.Label(janela, text="Letra da música (coloque o título da música na primeira linha):")
legenda["font"] = ("Calibri", "14")
legenda.pack()

# Frame para a barra de rolagem
frame = tk.Frame(janela)
frame.pack()

# Caixa de texto grande para inserção de dados
caixa_texto = tk.Text(frame, height=30, width=60)
caixa_texto.pack(side=tk.LEFT)

# Criando um menu contextual
menu_contextual = tk.Menu(janela, tearoff=0)
menu_contextual.add_command(label="Colar", command=colar_texto)
menu_contextual.add_command(label="Copiar", command=copiar_texto)

# Ligando o evento de clique do botão direito (abrir menu)
caixa_texto.bind("<Button-3>", abrir_menu)

# Barra de rolagem vertical
scrollbar_vertical = tk.Scrollbar(frame, command=caixa_texto.yview)
scrollbar_vertical.pack(side=tk.RIGHT, fill=tk.Y)
caixa_texto['yscrollcommand'] = scrollbar_vertical.set

# Botão para submeter o texto inserido
botao = tk.Button(janela, text="Criar slides", command=clique_botao, background="#555555", foreground="white", width=30)
botao["font"] = ("Calibri", "18")
botao.pack()

# Iniciar o loop principal da janela
janela.mainloop()